import streamlit as st
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain.chains.summarize import load_summarize_chain
from langchain.docstore.document import Document
import base64
import os
from dotenv import load_dotenv

# Optional: html templates
css = """
<style>
.justified { text-align: justify; }
</style>
"""
user_template = "<div style='background-color:#e0f7fa;padding:10px;border-radius:10px;margin:10px 0;'>User: {{MSG}}</div>"
bot_template = "<div style='background-color:#fff3e0;padding:10px;border-radius:10px;margin:10px 0;'>Bot: {{MSG}}</div>"

# Load API Key
load_dotenv()
key = os.getenv("OPENAI_API_KEY")
if not key:
    raise ValueError("❌ OPENAI_API_KEY is missing from the .env file!")
else:
    print("✅ API key loaded.")

# --------------------- File Parsing ------------------------- #
def extract_text_from_files(uploaded_files):
    documents = []
    for file in uploaded_files:
        if file.name.endswith(".pdf"):
            pdf_reader = PdfReader(file)
            for i, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text:
                    documents.append(Document(page_content=text, metadata={"source": file.name, "page": i + 1}))
        elif file.name.endswith(".docx"):
            docx = DocxDocument(file)
            text = "\n".join([para.text for para in docx.paragraphs])
            documents.append(Document(page_content=text, metadata={"source": file.name}))
        elif file.name.endswith(".pptx"):
            ppt = Presentation(file)
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text.strip() + "\n"
            documents.append(Document(page_content=text, metadata={"source": file.name}))
    return documents

# ------------------ LangChain Components -------------------- #
def split_documents(documents):
    text_splitter = CharacterTextSplitter(
        separator="\n", chunk_size=1000, chunk_overlap=200, length_function=len
    )
    return text_splitter.split_documents(documents)

def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings(openai_api_key=key)
    return FAISS.from_documents(text_chunks, embedding=embeddings)

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI(openai_api_key=key)
    memory = ConversationBufferMemory(
        memory_key='chat_history',
        return_messages=True,
        output_key='answer'
    )
    return ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(search_kwargs={"k": 3}),
        memory=memory,
        return_source_documents=True,
        output_key='answer'
    )

def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    answer = response['answer']
    source_docs = response.get('source_documents', [])
    sources = set(doc.metadata.get("source", "?") for doc in source_docs)
    page_info = f"<div style='color: gray; font-size: small;'>Sources: {', '.join(map(str, sources))}</div>"

    st.session_state.chat_history.append(("user", user_question))
    st.session_state.chat_history.append(("bot", answer))

    st.write(user_template.replace("{{MSG}}", user_question), unsafe_allow_html=True)
    st.write(bot_template.replace("{{MSG}}", answer + page_info), unsafe_allow_html=True)

def summarize_documents(vectorstore):
    llm = ChatOpenAI(openai_api_key=key, temperature=0)
    summarize_chain = load_summarize_chain(llm, chain_type="map_reduce")
    documents = list(vectorstore.docstore._dict.values())
    return summarize_chain.run(documents)

def download_chat_history():
    chat_lines = []
    for role, message in st.session_state.chat_history:
        prefix = "User: " if role == "user" else "Bot: "
        chat_lines.append(f"{prefix}{message}")
    history_text = "\n".join(chat_lines)
    st.download_button("💾 Download Chat History", data=history_text, file_name="chat_history.txt")

# ---------------------- File Preview ------------------------ #
def display_pdf(file):
    file.seek(0)
    base64_pdf = base64.b64encode(file.read()).decode('utf-8')
    st.markdown(f"""
        <iframe src="data:application/pdf;base64,{base64_pdf}"
                width="100%" height="700px" type="application/pdf"></iframe>
    """, unsafe_allow_html=True)

def display_docx(file):
    file.seek(0)
    doc = DocxDocument(file)
    content = ""
    for para in doc.paragraphs:
        content += f"<p>{para.text}</p>"
    st.markdown(content, unsafe_allow_html=True)

def display_pptx(file):
    file.seek(0)
    prs = Presentation(file)
    for idx, slide in enumerate(prs.slides):
        slide_text = f"<h4>Slide {idx + 1}</h4><ul>"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += f"<li>{shape.text.strip()}</li>"
        slide_text += "</ul><hr>"
        st.markdown(slide_text, unsafe_allow_html=True)

# ---------------------- Main App ---------------------------- #
def main():
    st.set_page_config(page_title="Chat with Files", page_icon="📚", layout="wide")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "vectorstore" not in st.session_state:
        st.session_state.vectorstore = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    st.markdown("<h1 style='text-align: center;'>📚 Chat with your Files (PDF, DOCX, PPTX)</h1>", unsafe_allow_html=True)

    col1, col2 = st.columns([1, 1.2], gap="large")

    with col1:
        with st.container(border=True, height=700):
            st.subheader("📂 Upload & Interact")
            uploaded_files = st.file_uploader(
                "Upload PDF, DOCX, or PPTX files", accept_multiple_files=True,
                type=["pdf", "docx", "pptx"]
            )

            if st.button("Process") and uploaded_files:
                with st.spinner("Processing files..."):
                    docs = extract_text_from_files(uploaded_files)
                    chunks = split_documents(docs)
                    vectorstore = get_vectorstore(chunks)
                    st.session_state.vectorstore = vectorstore
                    st.session_state.conversation = get_conversation_chain(vectorstore)
                    st.success("✅ Ready to chat!")

            question = st.text_input("Ask a question:")
            if question and st.session_state.conversation:
                handle_userinput(question)

            if st.button("Summarize") and st.session_state.vectorstore:
                with st.spinner("Generating summary..."):
                    summary = summarize_documents(st.session_state.vectorstore)
                st.markdown(f"<div class='justified'><strong>Summary:</strong><br>{summary}</div>", unsafe_allow_html=True)

            if st.session_state.chat_history:
                download_chat_history()

    with col2:
        with st.container(border=True, height=700):
            st.subheader("📄 File Preview")
            if uploaded_files:
                for file in uploaded_files:
                    with st.expander(file.name, expanded=False):
                        if file.name.endswith(".pdf"):
                            display_pdf(file)
                        elif file.name.endswith(".docx"):
                            display_docx(file)
                        elif file.name.endswith(".pptx"):
                            display_pptx(file)

if __name__ == '__main__':
    main()
